from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.styles import Alignment, Font, PatternFill
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / 'sample_uploads'
OUT_DIR.mkdir(exist_ok=True)

NAVY = PptRGBColor(15, 32, 56)
BLUE = PptRGBColor(0, 87, 164)
TEAL = PptRGBColor(20, 184, 166)
ORANGE = PptRGBColor(245, 130, 32)
GRAY = PptRGBColor(100, 116, 139)
LIGHT = PptRGBColor(241, 245, 249)
WHITE = PptRGBColor(255, 255, 255)


def add_textbox(slide, text, left, top, width, height, *, size=18, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = 'Aptos'
        run.font.size = PptPt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, title, value, caption, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(left), PptInches(top), PptInches(width), PptInches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = PptRGBColor(226, 232, 240)
    shape.line.width = PptPt(1)
    add_textbox(slide, title, left + 0.18, top + 0.18, width - 0.36, 0.25, size=9, color=GRAY, bold=True)
    add_textbox(slide, value, left + 0.18, top + 0.48, width - 0.36, 0.45, size=21, color=accent, bold=True)
    add_textbox(slide, caption, left + 0.18, top + 0.98, width - 0.36, 0.35, size=8, color=GRAY)


def add_header(slide, title, kicker='CONFIDENTIAL BUSINESS REVIEW'):
    ribbon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(13.333), PptInches(0.18))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = BLUE
    ribbon.line.fill.background()
    add_textbox(slide, kicker, 0.55, 0.28, 3.4, 0.25, size=7, color=GRAY, bold=True)
    add_textbox(slide, title, 0.55, 0.58, 8.8, 0.48, size=23, color=NAVY, bold=True)
    add_textbox(slide, 'Q1 2026', 11.35, 0.43, 1.35, 0.35, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tag = slide.shapes[-1]
    tag.fill.solid()
    tag.fill.fore_color.rgb = ORANGE
    tag.line.fill.background()


def add_footer(slide):
    add_textbox(slide, 'Acme Growth Co. | Board Review Template | Q1 2026', 0.55, 7.08, 7.0, 0.22, size=7, color=GRAY)
    add_textbox(slide, 'For internal planning discussion only', 9.3, 7.08, 3.2, 0.22, size=7, color=GRAY, align=PP_ALIGN.RIGHT)


def build_template_pptx():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    cover = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(13.333), PptInches(7.5))
    cover.fill.solid()
    cover.fill.fore_color.rgb = NAVY
    cover.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(5.95), PptInches(13.333), PptInches(1.55))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()
    add_textbox(slide, 'Q1 2026 Business Review', 0.75, 1.1, 7.5, 0.8, size=36, color=WHITE, bold=True)
    add_textbox(slide, 'Revenue momentum, customer health, and operating priorities', 0.78, 2.0, 7.4, 0.5, size=16, color=PptRGBColor(203, 213, 225))
    add_textbox(slide, 'Prepared for Board of Directors | April 2026', 0.78, 2.65, 5.2, 0.35, size=11, color=PptRGBColor(203, 213, 225))
    add_textbox(slide, 'Q1 closed with $12.4M revenue and 64% gross margin', 0.78, 6.25, 8.6, 0.42, size=18, color=WHITE, bold=True)
    add_textbox(slide, 'Template anchor: replace this slide with the latest quarter narrative.', 0.8, 6.78, 6.8, 0.25, size=8, color=PptRGBColor(203, 213, 225))

    slide = prs.slides.add_slide(blank)
    add_header(slide, 'Executive scorecard')
    add_card(slide, 0.65, 1.55, 2.85, 1.35, 'Revenue', '$12.4M', '+8% QoQ in Q1', BLUE)
    add_card(slide, 3.75, 1.55, 2.85, 1.35, 'Gross margin', '64%', '+2 pts vs Q4', TEAL)
    add_card(slide, 6.85, 1.55, 2.85, 1.35, 'Net revenue retention', '111%', 'Enterprise expansion led', ORANGE)
    add_card(slide, 9.95, 1.55, 2.85, 1.35, 'Pipeline coverage', '2.8x', 'Q2 target coverage', BLUE)
    add_textbox(slide, 'Q1 headline: efficient growth continued, but enterprise conversion timing remained the main constraint.', 0.75, 3.35, 11.7, 0.55, size=17, color=NAVY, bold=True)
    add_textbox(slide, 'Key observations\n- Enterprise ARR expanded 14% sequentially.\n- SMB churn remained elevated at 3.1% monthly.\n- Marketing CAC improved but sales cycle lengthened by 9 days.\n- Hiring stayed within the approved operating plan.', 0.85, 4.15, 5.65, 1.75, size=14, color=NAVY)
    add_textbox(slide, 'Management asks\n- Approve additional enterprise solution-engineering capacity.\n- Keep Q2 discretionary spend gated by pipeline conversion.\n- Revisit SMB packaging if churn remains above 2.5%.', 6.95, 4.15, 5.25, 1.75, size=14, color=NAVY)
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_header(slide, 'Revenue trend and forecast')
    chart_data = CategoryChartData()
    chart_data.categories = ['Oct', 'Nov', 'Dec', 'Jan', 'Feb', 'Mar']
    chart_data.add_series('Q1 Revenue', (3.5, 3.8, 4.0, 3.9, 4.1, 4.4))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, PptInches(0.75), PptInches(1.45), PptInches(7.25), PptInches(3.6), chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.chart_title.text_frame.text = 'Monthly revenue, Q1 baseline'
    add_textbox(slide, 'Q1 revenue closed at $12.4M, below the stretch case but ahead of the board-approved base plan.', 8.35, 1.55, 3.9, 0.78, size=16, color=NAVY, bold=True)
    add_textbox(slide, 'Forecast notes\n- March rebound came from two enterprise expansions.\n- New logo conversion was slower than planned.\n- Q2 forecast assumes no material price increase.\n- Upside depends on healthcare and fintech pipeline.', 8.35, 2.7, 3.8, 1.85, size=13, color=NAVY)
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_header(slide, 'Segment performance')
    table = slide.shapes.add_table(5, 5, PptInches(0.75), PptInches(1.45), PptInches(11.8), PptInches(2.65)).table
    values = [
        ['Segment', 'Q1 Revenue', 'QoQ Growth', 'Gross Margin', 'Commentary'],
        ['Enterprise', '$6.8M', '+14%', '71%', 'Expansion-led growth'],
        ['Mid-market', '$3.7M', '+6%', '63%', 'Stable demand'],
        ['SMB', '$1.9M', '-3%', '49%', 'Churn pressure'],
        ['Total', '$12.4M', '+8%', '64%', 'Mixed but improving'],
    ]
    for r, row in enumerate(values):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.margin_left = PptInches(0.05)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Aptos'
                    run.font.size = PptPt(9 if r else 8)
                    run.font.bold = r == 0 or c == 0
                    run.font.color.rgb = WHITE if r == 0 else NAVY
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
    add_textbox(slide, 'Q1 segment readout showed enterprise quality improving while SMB economics remained below threshold.', 0.85, 4.6, 8.6, 0.45, size=16, color=NAVY, bold=True)
    add_textbox(slide, 'Decision lens: prioritize segments with durable margin expansion and lower support burden.', 0.85, 5.15, 8.6, 0.35, size=12, color=GRAY)
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_header(slide, 'Risks, mitigations, and operating watchlist')
    add_textbox(slide, 'Q1 risk posture remained manageable, with customer concentration and SMB churn requiring continued attention.', 0.75, 1.35, 10.8, 0.5, size=17, color=NAVY, bold=True)
    add_textbox(slide, 'Risk register\n- Enterprise deals are concentrated in 6 named accounts.\n- SMB churn may offset self-serve acquisition gains.\n- Cloud infrastructure cost could rise with AI usage.\n- Hiring delays may constrain implementation capacity.', 0.85, 2.15, 5.6, 1.7, size=14, color=NAVY)
    add_textbox(slide, 'Mitigation plan\n- Expand executive sponsor coverage for top accounts.\n- Launch SMB onboarding redesign before end of Q2.\n- Add usage-based cost guardrails to AI workloads.\n- Pull forward two implementation hires if bookings convert.', 6.85, 2.15, 5.5, 1.7, size=14, color=NAVY)
    add_textbox(slide, 'Board discussion: Should management continue trading SMB growth for enterprise focus in Q2?', 0.9, 5.35, 11.0, 0.48, size=16, color=BLUE, bold=True)
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_header(slide, 'Q2 priorities and decision requests')
    add_textbox(slide, 'Q1-to-Q2 operating bridge', 0.75, 1.35, 4.4, 0.45, size=18, color=NAVY, bold=True)
    add_textbox(slide, '1. Convert enterprise late-stage pipeline\n2. Improve SMB onboarding and retention\n3. Protect margin while scaling AI features\n4. Maintain hiring discipline against revenue proof points', 0.9, 2.0, 5.15, 2.1, size=16, color=NAVY)
    add_textbox(slide, 'Q1 decision requests\n- Approve $420K incremental enterprise GTM capacity.\n- Hold Q2 brand spend until pipeline conversion improves.\n- Confirm risk appetite for AI feature infrastructure cost.', 6.85, 1.7, 5.3, 1.55, size=15, color=NAVY)
    add_textbox(slide, 'Expected Q2 board update: conversion quality, retention progress, and infrastructure cost envelope.', 6.85, 4.1, 5.0, 0.8, size=16, color=BLUE, bold=True)
    add_footer(slide)

    prs.save(OUT_DIR / 'template.pptx')


def style_doc(document: Document):
    styles = document.styles
    styles['Normal'].font.name = 'Aptos'
    styles['Normal'].font.size = Pt(10.5)
    for style_name in ['Title', 'Heading 1', 'Heading 2']:
        style = styles[style_name]
        style.font.name = 'Aptos Display'
        style.font.color.rgb = RGBColor(15, 32, 56)


def add_doc_table(document: Document, headers: list[str], rows: list[list[str]]):
    table = document.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = 'Table Grid'
    for idx, header in enumerate(headers):
        cell = table.rows[0].cells[idx]
        cell.text = header
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
    for row in rows:
        cells = table.add_row().cells
        for idx, value in enumerate(row):
            cells[idx].text = value
    document.add_paragraph()


def build_content_docx():
    document = Document()
    style_doc(document)
    section = document.sections[0]
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    title = document.add_paragraph()
    title.style = 'Title'
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.add_run('Q2 2026 Business Review Source Brief').bold = True
    subtitle = document.add_paragraph('Prepared for Board of Directors | July 2026 | Source content for Prompt-to-PPT transformation test')
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

    document.add_heading('Executive Summary', level=1)
    document.add_paragraph('Q2 closed materially ahead of the Q1 baseline. Revenue reached $15.8M, representing 27% quarter-over-quarter growth. Gross margin expanded to 72% as enterprise mix improved and cloud optimization efforts reduced unit cost.')
    document.add_paragraph('The management narrative should shift from Q1 efficiency recovery to Q2 durable enterprise acceleration. The board should see a clear upgrade in revenue quality, customer expansion, and pipeline confidence.')
    document.add_paragraph('Primary decision request: approve a targeted $650K investment in enterprise implementation capacity and customer success coverage while keeping brand spend gated by conversion evidence.')

    document.add_heading('Q2 Scorecard', level=1)
    add_doc_table(
        document,
        ['Metric', 'Q1 Baseline', 'Q2 Actual', 'Management Read'],
        [
            ['Revenue', '$12.4M', '$15.8M', 'Beat base plan by $1.1M and grew 27% QoQ'],
            ['Gross Margin', '64%', '72%', 'Enterprise mix and AI cost controls improved margin'],
            ['Net Revenue Retention', '111%', '118%', 'Expansion motion strengthened in top accounts'],
            ['Pipeline Coverage', '2.8x', '3.6x', 'Q3 coverage is healthier than Q2 entry point'],
            ['SMB Monthly Churn', '3.1%', '2.4%', 'Onboarding redesign is beginning to reduce churn'],
        ],
    )

    document.add_heading('Revenue Trend and Forecast', level=1)
    document.add_paragraph('Monthly revenue increased from $4.8M in April to $5.6M in June. The quarter was not dependent on one large transaction: the top three expansion deals represented 31% of incremental ARR versus 46% in Q1.')
    document.add_paragraph('Q3 forecast is $17.2M base case and $18.4M upside case. Upside depends on healthcare pipeline conversion and one strategic fintech renewal. Downside risk is implementation capacity, not demand generation.')

    document.add_heading('Segment Performance', level=1)
    add_doc_table(
        document,
        ['Segment', 'Q2 Revenue', 'QoQ Growth', 'Gross Margin', 'Board Message'],
        [
            ['Enterprise', '$8.9M', '+31%', '76%', 'Best-quality growth engine; capacity is the constraint'],
            ['Mid-market', '$4.4M', '+19%', '68%', 'Stable growth with improved sales efficiency'],
            ['SMB', '$2.5M', '+32%', '56%', 'Churn improved but support burden remains watch item'],
            ['Total', '$15.8M', '+27%', '72%', 'Mix shift supports both growth and margin expansion'],
        ],
    )

    document.add_heading('Customer and Pipeline Health', level=1)
    document.add_paragraph('Logo pipeline expanded to $57M qualified value, up from $42M at the end of Q1. Enterprise stage-3 conversion improved from 29% to 37% after solution-engineering coverage increased.')
    document.add_paragraph('Customer health is trending positive. Net revenue retention improved to 118%, gross logo retention reached 94%, and the number of accounts above $250K ARR increased from 18 to 24.')

    document.add_heading('Risks and Mitigations', level=1)
    document.add_paragraph('Implementation capacity is now the largest operating risk. If late-stage healthcare and fintech opportunities convert, current delivery staffing may create onboarding delays in September.')
    document.add_paragraph('SMB churn improved but remains sensitive to product activation time. Management will continue onboarding experiments and will pause additional self-serve acquisition spend if churn rises above 2.6%.')
    document.add_paragraph('AI infrastructure cost is within the approved envelope after model-routing optimization. The team recommends keeping usage guardrails active until customer-level profitability dashboards are complete.')

    document.add_heading('Q3 Priorities and Decision Requests', level=1)
    document.add_paragraph('Priority 1: convert enterprise late-stage pipeline without stretching implementation quality. Priority 2: sustain gross margin above 70%. Priority 3: reduce SMB activation time below 6 days. Priority 4: operationalize account health scoring for customer success.')
    document.add_paragraph('Board approval requested: $650K for enterprise implementation and customer success capacity, conditional release of $300K brand spend after two consecutive months of stage-3 conversion above 35%, and approval to expand AI infrastructure budget by up to $180K if tied to signed enterprise deployments.')

    document.save(OUT_DIR / 'content.docx')


def style_sheet(ws):
    header_fill = PatternFill('solid', fgColor='0F2038')
    for row in ws.iter_rows(min_row=1, max_row=1):
        for cell in row:
            cell.fill = header_fill
            cell.font = Font(color='FFFFFF', bold=True)
            cell.alignment = Alignment(horizontal='center')
    for column in ws.columns:
        width = max(len(str(cell.value)) if cell.value is not None else 0 for cell in column) + 3
        ws.column_dimensions[column[0].column_letter].width = min(max(width, 12), 28)


def build_content_xlsx():
    wb = Workbook()
    ws = wb.active
    ws.title = 'KPI Summary'
    ws.append(['Metric', 'Q1 Baseline', 'Q2 Actual', 'Delta', 'Target'])
    rows = [
        ['Revenue ($M)', 12.4, 15.8, 3.4, 14.7],
        ['Gross Margin (%)', 64, 72, 8, 70],
        ['NRR (%)', 111, 118, 7, 115],
        ['Pipeline Coverage (x)', 2.8, 3.6, 0.8, 3.2],
        ['SMB Churn (%)', 3.1, 2.4, -0.7, 2.6],
    ]
    for row in rows:
        ws.append(row)
    style_sheet(ws)
    chart = BarChart()
    chart.title = 'Q1 vs Q2 KPI Snapshot'
    chart.y_axis.title = 'Value'
    data = Reference(ws, min_col=2, max_col=3, min_row=1, max_row=6)
    cats = Reference(ws, min_col=1, min_row=2, max_row=6)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    chart.height = 8
    chart.width = 14
    ws.add_chart(chart, 'G2')

    ws = wb.create_sheet('Monthly Revenue')
    ws.append(['Month', 'Revenue $M', 'Enterprise $M', 'Mid-market $M', 'SMB $M'])
    for row in [
        ['Jan', 3.9, 2.1, 1.2, 0.6],
        ['Feb', 4.1, 2.2, 1.3, 0.6],
        ['Mar', 4.4, 2.5, 1.2, 0.7],
        ['Apr', 4.8, 2.7, 1.3, 0.8],
        ['May', 5.4, 3.1, 1.5, 0.8],
        ['Jun', 5.6, 3.1, 1.6, 0.9],
    ]:
        ws.append(row)
    style_sheet(ws)
    line = LineChart()
    line.title = 'Monthly Revenue Trend'
    line.y_axis.title = '$M'
    data = Reference(ws, min_col=2, max_col=5, min_row=1, max_row=7)
    cats = Reference(ws, min_col=1, min_row=2, max_row=7)
    line.add_data(data, titles_from_data=True)
    line.set_categories(cats)
    line.height = 8
    line.width = 14
    ws.add_chart(line, 'G2')

    ws = wb.create_sheet('Pipeline Forecast')
    ws.append(['Segment', 'Qualified Pipeline $M', 'Stage 3 Conversion %', 'Q3 Base Forecast $M', 'Q3 Upside $M'])
    for row in [
        ['Enterprise', 36.0, 37, 10.1, 11.2],
        ['Mid-market', 14.0, 31, 4.6, 4.9],
        ['SMB', 7.0, 18, 2.5, 2.3],
        ['Total', 57.0, 33, 17.2, 18.4],
    ]:
        ws.append(row)
    style_sheet(ws)

    ws = wb.create_sheet('Customer Health')
    ws.append(['Metric', 'Q1', 'Q2', 'Comment'])
    for row in [
        ['Accounts > $250K ARR', 18, 24, 'Enterprise expansion momentum'],
        ['Gross Logo Retention %', 91, 94, 'Improved onboarding and support response'],
        ['Median Activation Days', 8, 6, 'SMB onboarding redesign impact'],
        ['At-risk Enterprise Accounts', 7, 4, 'Executive sponsor program reduced risk'],
    ]:
        ws.append(row)
    style_sheet(ws)

    wb.properties.title = 'Q2 2026 Business Review Workbook'
    wb.save(OUT_DIR / 'content.xlsx')


if __name__ == '__main__':
    build_template_pptx()
    build_content_docx()
    build_content_xlsx()
    print(f'Wrote sample upload files to {OUT_DIR}')
