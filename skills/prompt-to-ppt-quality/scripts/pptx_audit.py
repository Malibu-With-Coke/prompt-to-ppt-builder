#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any


EMU_PER_INCH = 914400
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}


def emu_to_inches(value: Any) -> float:
    try:
        return round(int(value) / EMU_PER_INCH, 3)
    except Exception:
        return 0.0


def compact_text(value: Any, limit: int = 500) -> str:
    text = " ".join(str(value or "").split())
    if len(text) <= limit:
        return text
    return text[: limit - 1].rstrip() + "..."


def audit_pptx(path: Path) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ModuleNotFoundError:
        return audit_pptx_ooxml(path)

    presentation = Presentation(str(path))
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    slides = []
    fonts: set[str] = set()
    colors: set[str] = set()
    warnings: list[dict[str, Any]] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        text_shapes = []
        for shape in slide.shapes:
            shape_summary = summarize_shape(shape)
            if not shape_summary:
                continue
            text_shapes.append(shape_summary)
            collect_style(shape, fonts, colors)
            warnings.extend(shape_warnings(slide_index, shape_summary, slide_width, slide_height))

        slides.append(
            {
                "index": slide_index,
                "layoutName": slide.slide_layout.name if slide.slide_layout else "",
                "textShapeCount": len(text_shapes),
                "textShapes": text_shapes,
            }
        )

    return {
        "path": str(path),
        "engine": "python-pptx",
        "slideCount": len(presentation.slides),
        "slideWidth": slide_width,
        "slideHeight": slide_height,
        "slideSizeInches": {
            "width": emu_to_inches(slide_width),
            "height": emu_to_inches(slide_height),
        },
        "slides": slides,
        "fonts": sorted(fonts),
        "colors": sorted(colors)[:24],
        "warnings": warnings,
        "warningCount": len(warnings),
    }


def audit_pptx_ooxml(path: Path) -> dict[str, Any]:
    with zipfile.ZipFile(path) as package:
        presentation_xml = ET.fromstring(package.read("ppt/presentation.xml"))
        size = presentation_xml.find("p:sldSz", NS)
        slide_width = int(size.attrib.get("cx", 0)) if size is not None else 0
        slide_height = int(size.attrib.get("cy", 0)) if size is not None else 0
        slide_paths = get_slide_paths(package, presentation_xml)

        slides = []
        fonts: set[str] = set()
        colors: set[str] = set()
        warnings: list[dict[str, Any]] = []

        for slide_index, slide_path in enumerate(slide_paths, start=1):
            slide_xml = ET.fromstring(package.read(slide_path))
            text_shapes = []
            for shape in slide_xml.findall(".//p:sp", NS):
                shape_summary = summarize_ooxml_shape(shape)
                if not shape_summary:
                    continue
                text_shapes.append(shape_summary)
                collect_ooxml_style(shape, fonts, colors)
                warnings.extend(shape_warnings(slide_index, shape_summary, slide_width, slide_height))

            slides.append(
                {
                    "index": slide_index,
                    "layoutName": "",
                    "textShapeCount": len(text_shapes),
                    "textShapes": text_shapes,
                }
            )

    return {
        "path": str(path),
        "engine": "ooxml-fallback",
        "slideCount": len(slides),
        "slideWidth": slide_width,
        "slideHeight": slide_height,
        "slideSizeInches": {
            "width": emu_to_inches(slide_width),
            "height": emu_to_inches(slide_height),
        },
        "slides": slides,
        "fonts": sorted(fonts),
        "colors": sorted(colors)[:24],
        "warnings": warnings,
        "warningCount": len(warnings),
    }


def get_slide_paths(package: zipfile.ZipFile, presentation_xml: ET.Element) -> list[str]:
    rels_xml = ET.fromstring(package.read("ppt/_rels/presentation.xml.rels"))
    rel_by_id = {
        rel.attrib.get("Id"): rel.attrib.get("Target")
        for rel in rels_xml.findall("rel:Relationship", NS)
    }
    paths = []
    for slide_id in presentation_xml.findall(".//p:sldId", NS):
        rel_id = slide_id.attrib.get(f"{{{NS['r']}}}id")
        target = rel_by_id.get(rel_id)
        if not target:
            continue
        if target.startswith("/"):
            paths.append(target.lstrip("/"))
        elif target.startswith("ppt/"):
            paths.append(target)
        else:
            paths.append(f"ppt/{target}")
    return paths


def summarize_ooxml_shape(shape: ET.Element) -> dict[str, Any] | None:
    text = compact_text(" ".join(node.text or "" for node in shape.findall(".//a:t", NS)))
    if not text:
        return None

    c_nv_pr = shape.find(".//p:cNvPr", NS)
    shape_id = int(c_nv_pr.attrib.get("id", 0)) if c_nv_pr is not None else 0
    name = c_nv_pr.attrib.get("name", "") if c_nv_pr is not None else ""

    placeholder = shape.find(".//p:ph", NS)
    placeholder_type = placeholder.attrib.get("type") if placeholder is not None else None

    xfrm = shape.find(".//a:xfrm", NS)
    off = xfrm.find("a:off", NS) if xfrm is not None else None
    ext = xfrm.find("a:ext", NS) if xfrm is not None else None
    left = int(off.attrib.get("x", 0)) if off is not None else 0
    top = int(off.attrib.get("y", 0)) if off is not None else 0
    width = int(ext.attrib.get("cx", 0)) if ext is not None else 0
    height = int(ext.attrib.get("cy", 0)) if ext is not None else 0
    has_bounds = width > 0 and height > 0
    area_inches = max(emu_to_inches(width) * emu_to_inches(height), 0.01) if has_bounds else 0

    return {
        "shapeId": shape_id,
        "name": name,
        "placeholderType": placeholder_type,
        "text": text,
        "charCount": len(text),
        "position": {
            "left": left,
            "top": top,
            "width": width,
            "height": height,
            "leftIn": emu_to_inches(left),
            "topIn": emu_to_inches(top),
            "widthIn": emu_to_inches(width),
            "heightIn": emu_to_inches(height),
        },
        "hasBounds": has_bounds,
        "densityCharsPerSqIn": round(len(text) / area_inches, 1) if has_bounds else 0,
    }


def collect_ooxml_style(shape: ET.Element, fonts: set[str], colors: set[str]) -> None:
    for latin in shape.findall(".//a:latin", NS):
        typeface = latin.attrib.get("typeface")
        if typeface:
            fonts.add(typeface)
    for color in shape.findall(".//a:srgbClr", NS):
        value = color.attrib.get("val")
        if value:
            colors.add(value)


def summarize_shape(shape: Any) -> dict[str, Any] | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    text = compact_text(getattr(shape, "text", ""))
    if not text:
        return None

    placeholder_type = None
    if getattr(shape, "is_placeholder", False):
        placeholder_type = str(shape.placeholder_format.type)

    left = int(shape.left)
    top = int(shape.top)
    width = int(shape.width)
    height = int(shape.height)
    has_bounds = width > 0 and height > 0
    area_inches = max(emu_to_inches(width) * emu_to_inches(height), 0.01) if has_bounds else 0

    return {
        "shapeId": int(shape.shape_id),
        "name": shape.name,
        "placeholderType": placeholder_type,
        "text": text,
        "charCount": len(text),
        "position": {
            "left": left,
            "top": top,
            "width": width,
            "height": height,
            "leftIn": emu_to_inches(left),
            "topIn": emu_to_inches(top),
            "widthIn": emu_to_inches(width),
            "heightIn": emu_to_inches(height),
        },
        "hasBounds": has_bounds,
        "densityCharsPerSqIn": round(len(text) / area_inches, 1) if has_bounds else 0,
    }


def collect_style(shape: Any, fonts: set[str], colors: set[str]) -> None:
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.name:
                fonts.add(run.font.name)
            rgb = getattr(getattr(run.font, "color", None), "rgb", None)
            if rgb is not None:
                colors.add(str(rgb))


def shape_warnings(slide_index: int, shape: dict[str, Any], slide_width: int, slide_height: int) -> list[dict[str, Any]]:
    warnings = []
    pos = shape["position"]
    if pos["left"] < 0 or pos["top"] < 0 or pos["left"] + pos["width"] > slide_width or pos["top"] + pos["height"] > slide_height:
        warnings.append(
            {
                "type": "out_of_bounds",
                "slideIndex": slide_index,
                "shapeId": shape["shapeId"],
                "message": "Text shape extends beyond the slide canvas.",
            }
        )

    density = shape["densityCharsPerSqIn"]
    if shape.get("hasBounds") and density > 120:
        warnings.append(
            {
                "type": "high_text_density",
                "slideIndex": slide_index,
                "shapeId": shape["shapeId"],
                "densityCharsPerSqIn": density,
                "message": "Text may clip, wrap badly, or become too small after replacement.",
            }
        )

    if shape["charCount"] > 700:
        warnings.append(
            {
                "type": "long_text",
                "slideIndex": slide_index,
                "shapeId": shape["shapeId"],
                "charCount": shape["charCount"],
                "message": "Shape contains unusually long slide text.",
            }
        )

    return warnings


def main() -> int:
    parser = argparse.ArgumentParser(description="Audit PPTX text shapes for Prompt-to-PPT template extraction and output QA.")
    parser.add_argument("input", type=Path)
    parser.add_argument("--out", type=Path, help="Write JSON output to this path.")
    args = parser.parse_args()

    if not args.input.exists():
        raise SystemExit(f"Input not found: {args.input}")

    result = audit_pptx(args.input)
    payload = json.dumps(result, ensure_ascii=False, indent=2, default=str)
    if args.out:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(payload, encoding="utf-8")
    else:
        print(payload)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
