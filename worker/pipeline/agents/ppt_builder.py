from __future__ import annotations

import tempfile
from io import BytesIO
from pathlib import Path
from typing import Any

from utils.s3 import get_object_bytes


class PPTBuilder:
    def build(
        self,
        job: dict[str, Any],
        deck_transform: dict[str, Any],
    ) -> str:
        from pptx import Presentation

        template_bytes = get_object_bytes(job['templateS3Key'])
        presentation = Presentation(BytesIO(template_bytes))
        plan = deck_transform.get('transformPlan') or {}
        slide_plans = {int(slide['slideIndex']): slide for slide in plan.get('slides') or [] if slide.get('slideIndex')}

        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_plan = slide_plans.get(slide_index)
            if not slide_plan:
                continue
            replacements = {
                int(replacement['shapeId']): str(replacement.get('text') or '')
                for replacement in slide_plan.get('replacements') or []
                if replacement.get('shapeId') is not None
            }
            for shape in slide.shapes:
                if not getattr(shape, 'has_text_frame', False):
                    continue
                if shape.shape_id not in replacements:
                    continue
                self._replace_shape_text(shape, replacements[shape.shape_id])

            notes = slide.notes_slide.notes_text_frame
            notes.text = slide_plan.get('speakerNotes') or plan.get('strategy') or ''

        output_path = Path(tempfile.gettempdir()) / f'{job["jobId"]}-output.pptx'
        presentation.save(output_path)
        return str(output_path)

    def _replace_shape_text(self, shape: Any, text: str) -> None:
        from pptx.enum.text import MSO_AUTO_SIZE

        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        lines = text.splitlines() or ['']
        paragraphs = list(text_frame.paragraphs)

        for index, line in enumerate(lines):
            paragraph = paragraphs[index] if index < len(paragraphs) else text_frame.add_paragraph()
            self._replace_paragraph_text_preserving_style(paragraph, line)

        for paragraph in paragraphs[len(lines) :]:
            self._replace_paragraph_text_preserving_style(paragraph, '')

    def _replace_paragraph_text_preserving_style(self, paragraph: Any, text: str) -> None:
        runs = list(paragraph.runs)
        if not runs:
            paragraph.text = text
            return

        runs[0].text = text
        for run in runs[1:]:
            run.text = ''
