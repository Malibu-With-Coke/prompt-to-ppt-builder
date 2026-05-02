from __future__ import annotations

import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any


EMU_PER_INCH = 914400


class PPTValidationAgent:
    """Validate generated PPTX files with structural checks and optional render smoke tests."""

    def validate(self, job_id: str, output_path: str, deck_transform: dict[str, Any] | None = None) -> dict[str, Any]:
        path = Path(output_path)
        audit = self._audit_pptx(path)
        render_result = self._render_smoke_test(path)
        return {
            'jobId': job_id,
            'outputPath': str(path),
            'status': self._status(audit, render_result),
            'audit': audit,
            'render': render_result,
            'transformSummary': self._transform_summary(deck_transform or {}),
        }

    def _audit_pptx(self, path: Path) -> dict[str, Any]:
        from pptx import Presentation

        presentation = Presentation(str(path))
        slide_width = int(presentation.slide_width)
        slide_height = int(presentation.slide_height)
        warnings: list[dict[str, Any]] = []
        slides: list[dict[str, Any]] = []
        fonts: set[str] = set()
        colors: set[str] = set()

        for slide_index, slide in enumerate(presentation.slides, start=1):
            text_shapes = []
            for shape in slide.shapes:
                if not getattr(shape, 'has_text_frame', False):
                    continue
                text = self._compact_text(shape.text, 800)
                if not text:
                    continue
                shape_summary = self._summarize_shape(shape, text)
                text_shapes.append(shape_summary)
                warnings.extend(self._shape_warnings(slide_index, shape_summary, slide_width, slide_height))
                self._collect_style(shape, fonts, colors)
            slides.append(
                {
                    'index': slide_index,
                    'layoutName': slide.slide_layout.name if slide.slide_layout else '',
                    'textShapeCount': len(text_shapes),
                    'textShapes': text_shapes,
                }
            )

        return {
            'slideCount': len(presentation.slides),
            'slideWidth': slide_width,
            'slideHeight': slide_height,
            'slideSizeInches': {
                'width': self._emu_to_inches(slide_width),
                'height': self._emu_to_inches(slide_height),
            },
            'slides': slides,
            'fonts': sorted(fonts),
            'colors': sorted(colors)[:24],
            'warnings': warnings,
            'warningCount': len(warnings),
        }

    def _summarize_shape(self, shape: Any, text: str) -> dict[str, Any]:
        left = int(shape.left)
        top = int(shape.top)
        width = int(shape.width)
        height = int(shape.height)
        has_bounds = width > 0 and height > 0
        area_inches = max(self._emu_to_inches(width) * self._emu_to_inches(height), 0.01) if has_bounds else 0
        placeholder_type = None
        if getattr(shape, 'is_placeholder', False):
            placeholder_type = str(shape.placeholder_format.type)

        return {
            'shapeId': int(shape.shape_id),
            'name': shape.name,
            'placeholderType': placeholder_type,
            'text': text,
            'charCount': len(text),
            'hasBounds': has_bounds,
            'position': {
                'left': left,
                'top': top,
                'width': width,
                'height': height,
                'leftIn': self._emu_to_inches(left),
                'topIn': self._emu_to_inches(top),
                'widthIn': self._emu_to_inches(width),
                'heightIn': self._emu_to_inches(height),
            },
            'densityCharsPerSqIn': round(len(text) / area_inches, 1) if has_bounds else 0,
        }

    def _shape_warnings(self, slide_index: int, shape: dict[str, Any], slide_width: int, slide_height: int) -> list[dict[str, Any]]:
        warnings = []
        position = shape['position']
        if (
            position['left'] < 0
            or position['top'] < 0
            or position['left'] + position['width'] > slide_width
            or position['top'] + position['height'] > slide_height
        ):
            warnings.append(
                {
                    'type': 'out_of_bounds',
                    'slideIndex': slide_index,
                    'shapeId': shape['shapeId'],
                    'message': 'Text shape extends beyond the slide canvas.',
                }
            )
        if shape.get('hasBounds') and shape.get('densityCharsPerSqIn', 0) > 120:
            warnings.append(
                {
                    'type': 'high_text_density',
                    'slideIndex': slide_index,
                    'shapeId': shape['shapeId'],
                    'densityCharsPerSqIn': shape['densityCharsPerSqIn'],
                    'message': 'Text may clip, wrap badly, or become too small after replacement.',
                }
            )
        if shape.get('charCount', 0) > 700:
            warnings.append(
                {
                    'type': 'long_text',
                    'slideIndex': slide_index,
                    'shapeId': shape['shapeId'],
                    'charCount': shape['charCount'],
                    'message': 'Shape contains unusually long slide text.',
                }
            )
        return warnings

    def _collect_style(self, shape: Any, fonts: set[str], colors: set[str]) -> None:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.name:
                    fonts.add(run.font.name)
                color = getattr(getattr(run.font, 'color', None), 'rgb', None)
                if color is not None:
                    colors.add(str(color))

    def _render_smoke_test(self, path: Path) -> dict[str, Any]:
        office = shutil.which('soffice') or shutil.which('libreoffice')
        pdftoppm = shutil.which('pdftoppm')
        if not office or not pdftoppm:
            missing = []
            if not office:
                missing.append('soffice/libreoffice')
            if not pdftoppm:
                missing.append('pdftoppm')
            return {
                'status': 'skipped',
                'reason': f'Missing render tool(s): {", ".join(missing)}',
                'pageCount': 0,
                'previewFiles': [],
            }

        with tempfile.TemporaryDirectory(prefix='ppt-render-') as temp_dir:
            temp_path = Path(temp_dir)
            convert = subprocess.run(
                [office, '--headless', '--convert-to', 'pdf', '--outdir', str(temp_path), str(path)],
                capture_output=True,
                text=True,
                timeout=120,
                check=False,
            )
            pdf_path = temp_path / f'{path.stem}.pdf'
            if convert.returncode != 0 or not pdf_path.exists():
                return {
                    'status': 'failed',
                    'reason': self._compact_text(convert.stderr or convert.stdout or 'PPTX to PDF conversion failed.', 500),
                    'pageCount': 0,
                    'previewFiles': [],
                }

            preview_prefix = temp_path / 'slide'
            raster = subprocess.run(
                [pdftoppm, '-png', str(pdf_path), str(preview_prefix)],
                capture_output=True,
                text=True,
                timeout=120,
                check=False,
            )
            preview_files = sorted(str(file.name) for file in temp_path.glob('slide-*.png'))
            if raster.returncode != 0 or not preview_files:
                return {
                    'status': 'failed',
                    'reason': self._compact_text(raster.stderr or raster.stdout or 'PDF rasterization failed.', 500),
                    'pageCount': 0,
                    'previewFiles': [],
                }
            return {
                'status': 'passed',
                'reason': '',
                'pageCount': len(preview_files),
                'previewFiles': preview_files,
            }

    def _status(self, audit: dict[str, Any], render_result: dict[str, Any]) -> str:
        if render_result.get('status') == 'failed':
            return 'failed'
        if audit.get('warningCount', 0):
            return 'warning'
        return 'passed'

    def _transform_summary(self, deck_transform: dict[str, Any]) -> dict[str, Any]:
        plan = deck_transform.get('transformPlan') or {}
        slides = plan.get('slides') or []
        return {
            'deckTitle': plan.get('deckTitle'),
            'slideCount': len(slides),
            'replacementCount': sum(len(slide.get('replacements') or []) for slide in slides),
            'chartUpdateIntentCount': sum(len(slide.get('chartUpdates') or []) for slide in slides),
            'tableUpdateIntentCount': sum(len(slide.get('tableUpdates') or []) for slide in slides),
        }

    def _compact_text(self, value: Any, limit: int = 320) -> str:
        text = ' '.join(str(value or '').split())
        if len(text) <= limit:
            return text
        return text[: limit - 1].rstrip() + '...'

    def _emu_to_inches(self, value: Any) -> float:
        try:
            return round(int(value) / EMU_PER_INCH, 3)
        except Exception:
            return 0.0
