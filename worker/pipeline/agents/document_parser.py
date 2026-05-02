from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any

from utils.s3 import get_object_bytes


class DocumentParser:
    def __init__(self, bucket_name: str | None = None):
        self.bucket_name = bucket_name

    def parse(self, job: dict[str, Any]) -> dict[str, Any]:
        job_id = job['jobId']
        template_key = job['templateS3Key']
        content_keys = self._content_s3_keys(job)

        template_bytes = get_object_bytes(template_key, self.bucket_name)
        content_documents = []
        for index, content_key in enumerate(content_keys, start=1):
            content_bytes = get_object_bytes(content_key, self.bucket_name)
            content_documents.append(self._parse_content_document(index, content_key, content_bytes))

        return {
            'jobId': job_id,
            'templateRules': self._parse_template(template_bytes),
            'contentSummary': self._combine_content_documents(content_documents),
            'userOptions': job.get('options', {}),
            'sources': {
                'templateS3Key': template_key,
                'contentS3Key': content_keys[0] if content_keys else None,
                'contentS3Keys': content_keys,
            },
        }

    def _content_s3_keys(self, job: dict[str, Any]) -> list[str]:
        raw_keys = job.get('contentS3Keys')
        if isinstance(raw_keys, list):
            keys = [str(key).strip() for key in raw_keys if str(key).strip()]
            if keys:
                return keys
        legacy_key = job.get('contentS3Key')
        if legacy_key:
            return [str(legacy_key)]
        raise ValueError('Job does not contain contentS3Key or contentS3Keys.')

    def _parse_content_document(self, index: int, content_key: str, payload: bytes) -> dict[str, Any]:
        parsed = self._parse_content(content_key, payload)
        return {
            'sourceIndex': index,
            'sourceKey': content_key,
            'fileName': Path(content_key).name,
            **parsed,
        }

    def _combine_content_documents(self, documents: list[dict[str, Any]]) -> dict[str, Any]:
        if not documents:
            raise ValueError('No content documents were parsed.')

        if len(documents) == 1:
            single_document = dict(documents[0])
            single_document['sourceCount'] = 1
            single_document['documents'] = documents
            return single_document

        sections: list[dict[str, Any]] = []
        document_types = []
        for document in documents:
            document_types.append(document.get('documentType', 'unknown'))
            for section in document.get('sections') or []:
                sections.append(
                    {
                        **section,
                        'sourceIndex': document['sourceIndex'],
                        'sourceKey': document['sourceKey'],
                        'sourceTitle': document.get('title'),
                        'title': f'{document.get("title", document["fileName"])} - {section.get("title", "Section")}',
                    }
                )

        return {
            'title': f'Combined source content ({len(documents)} files)',
            'documentType': 'multi',
            'documentTypes': sorted(set(document_types)),
            'sourceCount': len(documents),
            'documents': documents,
            'sections': sections
            or [
                {
                    'title': 'Combined Overview',
                    'summary': 'The uploaded content files did not contain extractable sections.',
                    'dataType': 'text',
                    'bullets': [],
                }
            ],
            'workbookProfiles': [
                {
                    'sourceIndex': document['sourceIndex'],
                    'sourceTitle': document.get('title'),
                    **document.get('workbookProfile', {}),
                }
                for document in documents
                if document.get('workbookProfile')
            ],
            'documentProfiles': [
                {
                    'sourceIndex': document['sourceIndex'],
                    'sourceTitle': document.get('title'),
                    **document.get('documentProfile', {}),
                }
                for document in documents
                if document.get('documentProfile')
            ],
        }

    def _parse_content(self, content_key: str, payload: bytes) -> dict[str, Any]:
        lowered_key = content_key.lower()
        if lowered_key.endswith('.docx'):
            return self._parse_docx(payload)
        if lowered_key.endswith('.xlsx'):
            return self._parse_xlsx(payload)
        raise ValueError(f'Unsupported content file type: {content_key}')

    def _parse_template(self, payload: bytes) -> dict[str, Any]:
        from pptx import Presentation

        presentation = Presentation(BytesIO(payload))
        layout_summaries: list[dict[str, Any]] = []
        template_slides: list[dict[str, Any]] = []
        fonts: set[str] = set()
        colors: set[str] = set()
        max_bullets = 0

        for layout in presentation.slide_layouts:
            placeholders: list[str] = []
            for placeholder in layout.placeholders:
                placeholder_type = getattr(placeholder.placeholder_format, 'type', None)
                placeholders.append(str(placeholder_type) if placeholder_type is not None else 'UNKNOWN')
            layout_summaries.append(
                {
                    'name': layout.name or 'Unnamed Layout',
                    'placeholders': placeholders,
                }
            )

        for slide in presentation.slides:
            slide_shapes: list[dict[str, Any]] = []
            for shape in slide.shapes:
                shape_summary = self._summarize_template_shape(shape)
                if shape_summary:
                    slide_shapes.append(shape_summary)
                if not getattr(shape, 'has_text_frame', False):
                    continue
                text_frame = shape.text_frame
                bullet_count = sum(1 for paragraph in text_frame.paragraphs if paragraph.text.strip())
                max_bullets = max(max_bullets, bullet_count)
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            fonts.add(run.font.name)
                        color = getattr(getattr(run.font, 'color', None), 'rgb', None)
                        if color is not None:
                            colors.add(str(color))
            template_slides.append(
                {
                    'index': len(template_slides) + 1,
                    'layoutName': slide.slide_layout.name if slide.slide_layout else '',
                    'textShapes': slide_shapes,
                }
            )

        return {
            'slideCount': len(presentation.slides),
            'slideWidth': presentation.slide_width,
            'slideHeight': presentation.slide_height,
            'layouts': layout_summaries,
            'templateSlides': template_slides,
            'maxBullets': max_bullets or 4,
            'style': 'formal',
            'fonts': sorted(fonts),
            'colorTheme': sorted(colors)[:8],
        }

    def _summarize_template_shape(self, shape: Any) -> dict[str, Any] | None:
        if not getattr(shape, 'has_text_frame', False):
            return None
        text = ' '.join(shape.text.split())
        if not text:
            return None

        placeholder_type = None
        if getattr(shape, 'is_placeholder', False):
            placeholder_type = str(shape.placeholder_format.type)

        return {
            'shapeId': shape.shape_id,
            'name': shape.name,
            'placeholderType': placeholder_type,
            'text': text[:500],
            'position': {
                'left': int(shape.left),
                'top': int(shape.top),
                'width': int(shape.width),
                'height': int(shape.height),
            },
        }

    def _parse_docx(self, payload: bytes) -> dict[str, Any]:
        from docx import Document

        document = Document(BytesIO(payload))
        title = 'Uploaded document'
        sections: list[dict[str, Any]] = []
        current_section: dict[str, Any] | None = None
        heading_outline: list[dict[str, Any]] = []
        style_counts: dict[str, int] = {}
        paragraph_count = 0

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue

            style_name = paragraph.style.name if paragraph.style else ''
            style_counts[style_name or 'Unknown'] = style_counts.get(style_name or 'Unknown', 0) + 1
            paragraph_count += 1
            if style_name.startswith('Title') and title == 'Uploaded document':
                title = text
                continue

            if style_name.startswith('Heading'):
                if current_section:
                    sections.append(self._finalize_text_section(current_section))
                heading_outline.append(
                    {
                        'title': text,
                        'style': style_name,
                        'level': self._heading_level(style_name),
                    }
                )
                current_section = {'title': text, 'style': style_name, 'chunks': []}
                continue

            if current_section is None:
                current_section = {'title': 'Overview', 'style': 'Body', 'chunks': []}
                if title == 'Uploaded document':
                    title = text
            current_section['chunks'].append(text)

        if current_section:
            sections.append(self._finalize_text_section(current_section))

        table_profiles: list[dict[str, Any]] = []
        for index, table in enumerate(document.tables, start=1):
            preview_rows = []
            for row in table.rows[:8]:
                preview_rows.append([self._compact_text(cell.text, 160) for cell in row.cells])
            table_profile = {
                'title': f'Table {index}',
                'summary': f'DOCX table with {len(table.rows)} row(s) and {len(table.columns)} column(s).',
                'dataType': 'table',
                'bullets': [],
                'columns': preview_rows[0] if preview_rows else [],
                'sampleRows': preview_rows[1:],
                'tablePreview': preview_rows[:4],
                'rowCount': len(table.rows),
                'columnCount': len(table.columns),
            }
            table_profiles.append(table_profile)
            sections.append(table_profile)

        report_profile = {
            'paragraphCount': paragraph_count,
            'tableCount': len(document.tables),
            'headingOutline': heading_outline,
            'styleCounts': style_counts,
            'reportSignals': self._docx_report_signals(heading_outline, table_profiles),
        }

        return {
            'title': title,
            'documentType': 'docx',
            'sections': sections
            or [
                {
                    'title': 'Overview',
                    'summary': 'The document did not contain extractable sections.',
                    'dataType': 'text',
                    'bullets': [],
                }
            ],
            'documentProfile': report_profile,
        }

    def _finalize_text_section(self, section: dict[str, Any]) -> dict[str, Any]:
        chunks = [chunk for chunk in section.get('chunks', []) if chunk]
        summary = ' '.join(chunks[:2])[:320]
        return {
            'title': section['title'],
            'summary': summary or 'No summary available.',
            'dataType': 'text',
            'style': section.get('style'),
            'bullets': chunks[:5],
            'paragraphCount': len(chunks),
        }

    def _parse_xlsx(self, payload: bytes) -> dict[str, Any]:
        from openpyxl import load_workbook

        value_workbook = load_workbook(BytesIO(payload), data_only=True, read_only=True)
        formula_workbook = load_workbook(BytesIO(payload), data_only=False, read_only=False)
        title = formula_workbook.properties.title or formula_workbook.sheetnames[0] or 'Workbook'
        sections: list[dict[str, Any]] = []
        workbook_profile = {
            'sheetCount': len(formula_workbook.sheetnames),
            'sheets': formula_workbook.sheetnames,
            'hasFormulas': False,
            'tableCount': 0,
            'chartCount': 0,
        }

        for sheet in value_workbook.worksheets:
            formula_sheet = formula_workbook[sheet.title]
            rows: list[list[Any]] = []
            for raw_row in sheet.iter_rows(values_only=True):
                normalized_row = [self._normalize_cell_value(cell) for cell in raw_row]
                if any(cell not in (None, '') for cell in normalized_row):
                    rows.append(normalized_row)
                if len(rows) >= 12:
                    break

            formula_cells = self._collect_formula_cells(formula_sheet)
            tables = self._collect_xlsx_tables(formula_sheet)
            charts = self._collect_xlsx_charts(formula_sheet)
            workbook_profile['hasFormulas'] = bool(workbook_profile['hasFormulas'] or formula_cells)
            workbook_profile['tableCount'] += len(tables)
            workbook_profile['chartCount'] += len(charts)

            if not rows and not formula_cells:
                continue

            headers = self._infer_headers(rows)
            sample_rows = rows[1:8] if len(rows) > 1 else []
            numeric_columns = self._infer_numeric_columns(headers, sample_rows)
            data_type = 'chart' if numeric_columns else 'table'
            sections.append(
                {
                    'title': sheet.title,
                    'summary': self._summarize_sheet(headers, sample_rows, numeric_columns, formula_cells, tables, charts),
                    'dataType': data_type,
                    'columns': headers,
                    'sampleRows': sample_rows,
                    'numericColumns': numeric_columns,
                    'formulaCells': formula_cells,
                    'tables': tables,
                    'charts': charts,
                    'sheetStats': {
                        'maxRow': formula_sheet.max_row,
                        'maxColumn': formula_sheet.max_column,
                        'sampledRowCount': len(sample_rows),
                    },
                }
            )

        return {
            'title': title,
            'documentType': 'xlsx',
            'workbookProfile': workbook_profile,
            'sections': sections
            or [
                {
                    'title': 'Workbook Overview',
                    'summary': 'The workbook did not contain extractable rows.',
                    'dataType': 'table',
                    'columns': [],
                    'sampleRows': [],
                    'numericColumns': [],
                }
            ],
        }

    def _summarize_sheet(
        self,
        headers: list[str],
        sample_rows: list[list[Any]],
        numeric_columns: list[str],
        formula_cells: list[dict[str, str]],
        tables: list[dict[str, str]],
        charts: list[dict[str, Any]],
    ) -> str:
        header_preview = ', '.join(headers[:4]) if headers else 'no headers'
        row_count = len(sample_rows)
        details = [f'Sheet contains {row_count} sampled data rows. Headers include {header_preview}.']
        if numeric_columns:
            numeric_preview = ', '.join(numeric_columns[:3])
            details.append(f'Chart-friendly numeric columns include {numeric_preview}.')
        if formula_cells:
            details.append(f'{len(formula_cells)} formula cells were sampled.')
        if tables:
            details.append(f'{len(tables)} structured Excel table(s) were detected.')
        if charts:
            details.append(f'{len(charts)} embedded chart object(s) were detected.')
        return ' '.join(details)[:520]

    def _compact_text(self, value: Any, limit: int = 320) -> str:
        text = ' '.join(str(value or '').split())
        if len(text) <= limit:
            return text
        return text[: limit - 1].rstrip() + '...'

    def _heading_level(self, style_name: str) -> int | None:
        suffix = style_name.replace('Heading', '').strip()
        try:
            return int(suffix)
        except ValueError:
            return None

    def _docx_report_signals(self, headings: list[dict[str, Any]], tables: list[dict[str, Any]]) -> dict[str, Any]:
        heading_titles = [heading.get('title', '').lower() for heading in headings]
        return {
            'hasExecutiveSummary': any('executive' in title or 'summary' in title or '\uc694\uc57d' in title for title in heading_titles),
            'hasFinancialOrMetricTables': any(
                any(
                    keyword in ' '.join(str(cell).lower() for cell in table.get('columns') or [])
                    for keyword in ['revenue', 'margin', 'cost', '\ub9e4\ucd9c', '\ube44\uc6a9']
                )
                for table in tables
            ),
            'suggestedDeckUse': 'report' if len(headings) >= 3 or tables else 'narrative',
        }

    def _normalize_cell_value(self, value: Any) -> Any:
        if value is None:
            return None
        if isinstance(value, (int, float, bool)):
            return value
        return self._compact_text(value, 160)

    def _infer_headers(self, rows: list[list[Any]]) -> list[str]:
        if not rows:
            return []
        width = max(len(row) for row in rows[:3])
        header_row = rows[0]
        return [
            str(header_row[index]).strip() if index < len(header_row) and header_row[index] not in (None, '') else f'Column {index + 1}'
            for index in range(width)
        ]

    def _infer_numeric_columns(self, headers: list[str], sample_rows: list[list[Any]]) -> list[str]:
        numeric_columns = []
        for index, header in enumerate(headers):
            column_values = [
                row[index]
                for row in sample_rows
                if index < len(row) and isinstance(row[index], (int, float)) and not isinstance(row[index], bool)
            ]
            if len(column_values) >= 2:
                numeric_columns.append(header)
        return numeric_columns

    def _collect_formula_cells(self, sheet: Any, limit: int = 40) -> list[dict[str, str]]:
        formula_cells: list[dict[str, str]] = []
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and cell.value.startswith('='):
                    formula_cells.append({'cell': cell.coordinate, 'formula': cell.value})
                    if len(formula_cells) >= limit:
                        return formula_cells
        return formula_cells

    def _collect_xlsx_tables(self, sheet: Any) -> list[dict[str, str]]:
        tables = []
        for name, table in getattr(sheet, 'tables', {}).items():
            ref = getattr(table, 'ref', None) or str(table)
            tables.append({'name': str(name), 'ref': str(ref)})
        return tables

    def _collect_xlsx_charts(self, sheet: Any) -> list[dict[str, Any]]:
        charts = []
        for index, chart in enumerate(getattr(sheet, '_charts', []) or [], start=1):
            charts.append(
                {
                    'index': index,
                    'type': chart.__class__.__name__,
                    'title': self._compact_text(getattr(chart, 'title', None), 120) if getattr(chart, 'title', None) else '',
                }
            )
        return charts
